"""
ファイル処理ユーティリティ
Discord添付ファイルやリンクからコンテンツを抽出
"""
import aiohttp
import io
import re
from typing import Optional, Dict, List, Tuple
from PyPDF2 import PdfReader
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image
import base64


class FileProcessor:
    """ファイル処理クラス"""

    # サポートするファイル形式
    SUPPORTED_IMAGES = ['.png', '.jpg', '.jpeg', '.gif', '.webp']
    SUPPORTED_DOCS = ['.pdf', '.xlsx', '.xls', '.pptx', '.ppt', '.txt', '.md', '.csv']

    @staticmethod
    async def download_file(url: str) -> Optional[bytes]:
        """URLからファイルをダウンロード"""
        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url, timeout=aiohttp.ClientTimeout(total=30)) as response:
                    if response.status == 200:
                        return await response.read()
                    else:
                        print(f"Failed to download file: {response.status}")
                        return None
        except Exception as e:
            print(f"Error downloading file: {e}")
            return None

    @staticmethod
    def extract_urls(text: str) -> List[str]:
        """テキストからURLを抽出"""
        url_pattern = r'https?://[^\s<>"{}|\\^`\[\]]+'
        return re.findall(url_pattern, text)

    @staticmethod
    def is_google_drive_url(url: str) -> bool:
        """Google Drive URLかどうか判定"""
        return 'drive.google.com' in url or 'docs.google.com' in url

    @staticmethod
    async def process_pdf(file_bytes: bytes) -> str:
        """PDFファイルからテキストを抽出"""
        try:
            pdf_file = io.BytesIO(file_bytes)
            reader = PdfReader(pdf_file)

            text_content = []
            for i, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    text_content.append(f"--- Page {i} ---\n{text}")

            if not text_content:
                return "PDFファイルからテキストを抽出できませんでした（画像のみのPDFの可能性があります）"

            return "\n\n".join(text_content)
        except Exception as e:
            return f"PDF処理エラー: {str(e)}"

    @staticmethod
    async def process_excel(file_bytes: bytes) -> str:
        """Excelファイルからテキストを抽出"""
        try:
            excel_file = io.BytesIO(file_bytes)
            workbook = load_workbook(excel_file, data_only=True)

            text_content = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"=== シート: {sheet_name} ===\n")

                # 最大100行まで読み込み
                max_rows = min(sheet.max_row, 100)
                for row in sheet.iter_rows(min_row=1, max_row=max_rows, values_only=True):
                    row_data = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_data):  # 空行をスキップ
                        text_content.append(" | ".join(row_data))

                if sheet.max_row > 100:
                    text_content.append(f"\n... (残り {sheet.max_row - 100} 行省略)")

            return "\n".join(text_content)
        except Exception as e:
            return f"Excel処理エラー: {str(e)}"

    @staticmethod
    async def process_powerpoint(file_bytes: bytes) -> str:
        """PowerPointファイルからテキストを抽出"""
        try:
            ppt_file = io.BytesIO(file_bytes)
            presentation = Presentation(ppt_file)

            text_content = []
            for i, slide in enumerate(presentation.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                if slide_text:
                    text_content.append(f"--- スライド {i} ---\n" + "\n".join(slide_text))

            if not text_content:
                return "PowerPointファイルからテキストを抽出できませんでした"

            return "\n\n".join(text_content)
        except Exception as e:
            return f"PowerPoint処理エラー: {str(e)}"

    @staticmethod
    async def process_text(file_bytes: bytes) -> str:
        """テキストファイルを読み込み"""
        try:
            # UTF-8でデコード、失敗したらShift-JISを試す
            try:
                return file_bytes.decode('utf-8')
            except UnicodeDecodeError:
                return file_bytes.decode('shift-jis', errors='ignore')
        except Exception as e:
            return f"テキスト処理エラー: {str(e)}"

    @staticmethod
    async def process_image(file_bytes: bytes, filename: str) -> Tuple[str, str]:
        """
        画像ファイルを処理してBase64エンコード
        Returns: (media_type, base64_data)
        """
        try:
            # 画像形式を判定
            image = Image.open(io.BytesIO(file_bytes))
            format_lower = image.format.lower() if image.format else 'png'

            # メディアタイプのマッピング
            media_type_map = {
                'jpeg': 'image/jpeg',
                'jpg': 'image/jpeg',
                'png': 'image/png',
                'gif': 'image/gif',
                'webp': 'image/webp'
            }

            media_type = media_type_map.get(format_lower, 'image/png')

            # Base64エンコード
            base64_data = base64.b64encode(file_bytes).decode('utf-8')

            return media_type, base64_data
        except Exception as e:
            print(f"Image processing error: {e}")
            return None, None

    @classmethod
    async def process_attachment(cls, attachment) -> Dict:
        """
        Discord添付ファイルを処理
        Returns: {
            'type': 'text' | 'image',
            'content': str (textの場合) or {'media_type': str, 'data': str} (imageの場合),
            'filename': str,
            'error': str (optional)
        }
        """
        filename = attachment.filename.lower()
        result = {
            'filename': attachment.filename,
            'type': None,
            'content': None,
            'error': None
        }

        # ファイルをダウンロード
        file_bytes = await cls.download_file(attachment.url)
        if not file_bytes:
            result['error'] = "ファイルのダウンロードに失敗しました"
            return result

        # 画像ファイル
        if any(filename.endswith(ext) for ext in cls.SUPPORTED_IMAGES):
            media_type, base64_data = await cls.process_image(file_bytes, attachment.filename)
            if media_type and base64_data:
                result['type'] = 'image'
                result['content'] = {
                    'media_type': media_type,
                    'data': base64_data
                }
            else:
                result['error'] = "画像の処理に失敗しました"

        # PDFファイル
        elif filename.endswith('.pdf'):
            result['type'] = 'text'
            result['content'] = await cls.process_pdf(file_bytes)

        # Excelファイル
        elif filename.endswith(('.xlsx', '.xls')):
            result['type'] = 'text'
            result['content'] = await cls.process_excel(file_bytes)

        # PowerPointファイル
        elif filename.endswith(('.pptx', '.ppt')):
            result['type'] = 'text'
            result['content'] = await cls.process_powerpoint(file_bytes)

        # テキストファイル
        elif filename.endswith(('.txt', '.md', '.csv')):
            result['type'] = 'text'
            result['content'] = await cls.process_text(file_bytes)

        else:
            result['error'] = f"未対応のファイル形式: {filename}"

        return result
